"""PPT 解析入口：Presentation → Presentation 模型。

打开 pptx，读取尺寸，遍历 slides。
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation as PptxPresentation

from ppt_transfor.models.schema import Presentation
from ppt_transfor.parser.slide import parse_slide


def parse_presentation(path: str | Path, out_dir: Path | None = None) -> Presentation:
    """解析 PPT 文件为 Presentation 模型

    Args:
        path: pptx 文件路径
        out_dir: 输出根目录（如 out/），用于计算图片输出目录 out/media。
            传入时图片提取为独立文件（image_path 引用），
            None 时降级为 base64 内嵌（向后兼容/测试场景）。

    Returns:
        Presentation 模型
    """
    path = Path(path)
    prs = PptxPresentation(str(path))

    # 计算 media_dir：out_dir/media，图片提取为独立文件
    # 不传 out_dir 时降级为 base64 内嵌，保持向后兼容
    media_dir = None
    if out_dir is not None:
        media_dir = Path(out_dir) / "media"

    model = Presentation(
        source_file=path.name,
        slide_width=int(prs.slide_width),
        slide_height=int(prs.slide_height),
    )

    for idx, slide in enumerate(prs.slides):
        model.slides.append(parse_slide(slide, idx, prs, media_dir))

    return model
