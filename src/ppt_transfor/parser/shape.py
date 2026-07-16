"""形状解析器：通用属性提取 + 按类型分发。

通用属性：shape_id/name/shape_type/位置/旋转/填充/边框/阴影
类型分发：table/group/picture/connector/autoshape/text_box/placeholder
对 placeholder 形状，解析继承自 layout/master 的对齐、字号、颜色。
"""

from __future__ import annotations

from os.path import normpath
from pathlib import Path

from lxml import etree
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ppt_transfor.models.schema import Color, Fill, GradientStop, Line, Shadow, Shape
from ppt_transfor.utils.color import parse_color

# XML 命名空间
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_C = "http://schemas.openxmlformats.org/drawingml/2006/chart"


def _parse_gradient_from_xml(shape_element, prs=None) -> tuple[str | None, float | None, list[GradientStop]]:
    """从 <a:spPr> 下读取渐变填充的 stops、类型和角度。

    Returns:
        (gradient_type, gradient_angle, gradient_stops)
    """
    stops: list[GradientStop] = []
    gradient_type = None
    gradient_angle = None

    if shape_element is None:
        return gradient_type, gradient_angle, stops

    try:
        spPr = shape_element.find(f"{{{NS_P}}}spPr")
        if spPr is None:
            return gradient_type, gradient_angle, stops

        gradFill = spPr.find(f"{{{NS_A}}}gradFill")
        if gradFill is None:
            return gradient_type, gradient_angle, stops

        # 渐变类型与角度
        lin = gradFill.find(f"{{{NS_A}}}lin")
        if lin is not None:
            gradient_type = "linear"
            ang = lin.get("ang")
            if ang is not None:
                try:
                    gradient_angle = float(ang)
                except (ValueError, TypeError):
                    pass

        path = gradFill.find(f"{{{NS_A}}}path")
        if path is not None:
            gradient_type = path.get("path", "path")

        rect = gradFill.find(f"{{{NS_A}}}rect")
        if rect is not None:
            gradient_type = "rect"

        # stops
        gsLst = gradFill.find(f"{{{NS_A}}}gsLst")
        if gsLst is not None:
            for gs in gsLst:
                pos = gs.get("pos")
                if pos is None:
                    continue
                try:
                    position = int(pos) / 100000.0
                except (ValueError, TypeError):
                    continue

                color = None
                # srgbClr：优先直接子元素，其次兼容 <a:solidFill>/<a:srgbClr> 写法
                srgb = gs.find(f"{{{NS_A}}}srgbClr")
                if srgb is None:
                    solid_fill = gs.find(f"{{{NS_A}}}solidFill")
                    if solid_fill is not None:
                        srgb = solid_fill.find(f"{{{NS_A}}}srgbClr")
                if srgb is not None:
                    val = srgb.get("val")
                    if val:
                        color = Color(type="rgb", value=val)

                # schemeClr：固化为 RGB；同样兼容直接子元素或包在 solidFill 中
                if color is None:
                    scheme = gs.find(f"{{{NS_A}}}schemeClr")
                    if scheme is None:
                        solid_fill = gs.find(f"{{{NS_A}}}solidFill")
                        if solid_fill is not None:
                            scheme = solid_fill.find(f"{{{NS_A}}}schemeClr")
                    if scheme is not None:
                        val = scheme.get("val")
                        if val:
                            from ppt_transfor.utils.inheritance import _resolve_schemeclr_to_rgb

                            # 传入 scheme 元素以应用 lumOff/lumMod 等修饰符
                            rgb_value = _resolve_schemeclr_to_rgb(val, prs, scheme)
                            if rgb_value is not None:
                                color = Color(type="rgb", value=rgb_value)
                            else:
                                color = Color(type="theme", value=val)

                if color is not None:
                    stops.append(GradientStop(position=position, color=color))
    except Exception:
        pass

    return gradient_type, gradient_angle, stops


def _parse_fill(fill, shape_element=None, prs=None) -> Fill | None:
    """解析填充。

    优先读取 XML：若 <a:spPr> 下显式存在 <a:noFill/>，则视为透明填充。
    python-pptx 对 noFill 形状有时会报告 fill.type=BACKGROUND，
    直接按 XML 判断可避免透明文本框被误判为背景填充。

    None（继承）返回 None，避免往返时因 add_shape 默认填充行为差异产生噪声。
    """
    # 优先检查 XML：显式 noFill 表示透明
    if shape_element is not None:
        try:
            spPr = shape_element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}spPr")
            if spPr is not None:
                no_fill = spPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}noFill")
                if no_fill is not None:
                    model = Fill()
                    model.type = "none"
                    return model
        except Exception:
            pass

    try:
        fill_type = fill.type
    except Exception:
        return None

    if fill_type is None:
        return None  # 继承样式，不记录

    # fill_type 是 MSO_FILL 枚举
    type_name = fill_type.name if hasattr(fill_type, "name") else str(fill_type)

    # BACKGROUND 表示“跟随幻灯片背景”，需保留以正确回写
    if type_name == "BACKGROUND":
        model = Fill()
        model.type = "background"
        return model

    model = Fill()
    if type_name == "SOLID":
        try:
            color = parse_color(fill.fore_color, prs)
        except Exception:
            color = None
        if color is not None:
            model.type = "solid"
            model.color = color
        else:
            # 无法解析颜色时降级为无填充，避免渲染为默认黑色
            model.type = "none"
        return model
    elif type_name == "GRADIENT":
        model.type = "gradient"
        grad_type, grad_angle, grad_stops = _parse_gradient_from_xml(shape_element, prs)
        if grad_type is not None:
            model.gradient_type = grad_type
        if grad_angle is not None:
            model.gradient_angle = grad_angle
        if grad_stops:
            model.gradient_stops = grad_stops
    elif type_name == "PATTERNED":
        model.type = "pattern"
    elif type_name == "PICTURE":
        model.type = "picture"
    else:
        model.type = type_name.lower()

    return model


def _parse_line(line, prs=None) -> Line | None:
    """解析边框线条。width=0 视为无线条（add_shape 默认行为）。"""
    model = Line()

    has_value = False

    # 先获取 <a:ln> 原始状态，避免后续访问 line.color 时 python-pptx 自动添加 <a:solidFill/>
    # 污染 XML 导致 no_fill 检测失败
    ln_el = None
    orig_has_solid_fill = False
    orig_has_no_fill = False
    try:
        from pptx.oxml.ns import qn

        ln_el = line._ln
        if ln_el is not None:
            orig_has_solid_fill = ln_el.find(qn("a:solidFill")) is not None
            orig_has_no_fill = ln_el.find(qn("a:noFill")) is not None
    except Exception:
        pass

    try:
        if line.width is not None:
            width_val = int(line.width)
            # width=0 视为无线条，不记录
            if width_val > 0:
                model.width = width_val
                has_value = True
    except Exception:
        pass

    try:
        color = parse_color(line.color, prs)
        if color is not None:
            model.color = color
            has_value = True
    except Exception:
        pass

    try:
        dash = line.dash_style
        if dash is not None:
            model.dash = dash.name if hasattr(dash, "name") else str(dash)
            has_value = True
    except Exception:
        pass

    # 检测无填充线条：<a:ln> 有 width 但无 solidFill（无论是否有 noFill）
    # 渲染时需显式设置 noFill，避免继承 p:style 或主题默认色（蓝色）
    # 使用最初获取的原始状态，避免 line.color 访问污染
    if ln_el is not None and model.width is not None:
        if not orig_has_solid_fill:
            # 无 solidFill（可能有 noFill 也可能没有）→ 标记为 no_fill
            model.no_fill = True
            has_value = True

    # 箭头线端点：从 <a:ln> 的 <a:headEnd>/<a:tailEnd> 读取 type 属性
    try:
        if ln_el is not None:
            head_end = ln_el.find(qn("a:headEnd"))
            if head_end is not None:
                head_type = head_end.get("type")
                if head_type:
                    model.head_arrow_type = head_type
                    has_value = True
            tail_end = ln_el.find(qn("a:tailEnd"))
            if tail_end is not None:
                tail_type = tail_end.get("type")
                if tail_type:
                    model.tail_arrow_type = tail_type
                    has_value = True
    except Exception:
        pass

    return model if has_value else None


def _parse_shadow(shape) -> Shadow | None:
    """解析阴影。

    inherit=False（显式无阴影）与 None（不记录）统一，
    避免往返时因 add_shape 默认继承阴影产生噪声。
    """
    try:
        shadow = shape.shadow
        if shadow is not None:
            # inherit=False 表示显式关闭阴影，与不记录统一
            if not shadow.inherit:
                return None
            return Shadow(enabled=True)
    except Exception:
        pass
    return None


def _parse_common_props(shape, prs=None) -> dict:
    """解析通用属性（位置/旋转/填充/边框/阴影）"""
    props = {}

    # 位置与尺寸（EMU 整数）
    try:
        props["left"] = int(shape.left) if shape.left is not None else None
    except Exception:
        props["left"] = None
    try:
        props["top"] = int(shape.top) if shape.top is not None else None
    except Exception:
        props["top"] = None
    try:
        props["width"] = int(shape.width) if shape.width is not None else None
    except Exception:
        props["width"] = None
    try:
        props["height"] = int(shape.height) if shape.height is not None else None
    except Exception:
        props["height"] = None

    # 旋转角度（度，浮点）
    try:
        props["rotation"] = float(shape.rotation or 0.0)
    except Exception:
        props["rotation"] = 0.0

    # 填充
    try:
        fill = _parse_fill(shape.fill, getattr(shape, "_element", None), prs)
        if fill is not None:
            props["fill"] = fill
    except Exception:
        pass

    # 边框
    try:
        line = _parse_line(shape.line, prs)
        if line is not None:
            props["line"] = line
    except Exception:
        pass

    # 阴影
    shadow = _parse_shadow(shape)
    if shadow is not None:
        props["shadow"] = shadow

    return props


def _is_placeholder(shape) -> bool:
    """判断形状是否为 placeholder"""
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER
    except Exception:
        return False


def _parse_chart_part(shape, slide) -> str | None:
    """从 slide part 的 relationships 中定位 chart 对应的 part 路径。

    返回标准化路径，例如 "ppt/charts/chart1.xml"。
    """
    if slide is None:
        return None

    try:
        graphic_data = shape._element.find(f".//{{{NS_A}}}graphicData")
        if graphic_data is None:
            return None

        chart_el = graphic_data.find(f"{{{NS_C}}}chart")
        if chart_el is None:
            return None

        chart_rid = chart_el.get(f"{{{NS_R}}}id")
        if not chart_rid:
            return None

        rel = slide.part.rels.get(chart_rid)
        if rel is None:
            return None

        target = rel.target_ref
        if not target:
            return None

        # target_ref 相对于 slide 文件（ppt/slides/）
        part_path = Path("ppt/slides") / target
        return normpath(str(part_path)).replace("\\", "/")
    except Exception:
        return None


def parse_shape(shape, slide=None, prs=None) -> Shape:
    """解析单个形状，按类型分发到对应解析器。

    Args:
        shape: python-pptx Shape 对象
        slide: 所属幻灯片（用于 placeholder 继承解析）
        prs: 所属 Presentation（用于主题色固化）
    """
    # 基础属性
    model = Shape(
        shape_id=str(getattr(shape, "shape_id", "") or ""),
        name=getattr(shape, "name", None),
    )

    # 检查是否有 <p:style> 元素（主题样式引用）
    # add_shape 会自动创建 p:style，但原始 PPT 可能没有
    try:
        p_style = shape._element.find(f"{{{NS_P}}}style")
        if p_style is not None:
            model.has_style = True
            # 记录 p:style 的完整 XML，用于回写
            # add_textbox 不创建 p:style，需要从原始 PPT 复制
            model.style_xml = etree.tostring(p_style, encoding="unicode")
        else:
            model.has_style = False
    except Exception:
        model.has_style = False

    # 通用属性
    props = _parse_common_props(shape, prs)
    for k, v in props.items():
        setattr(model, k, v)

    # 对 placeholder 形状，预先解析继承属性（对齐/字号/颜色/字体名），传入 prs 以固化主题色
    inherited_props = None
    if _is_placeholder(shape) and slide is not None:
        from ppt_transfor.utils.inheritance import resolve_placeholder_props
        inherited_props = resolve_placeholder_props(shape, slide, prs)
        # 不对齐方式做 placeholder 类型推断（如假设 TITLE 默认居中）。
        # OOXML 规范中 <a:pPr> 缺失 algn 属性时默认左对齐，
        # 若继承链（slide→layout→master）均未显式设置对齐，
        # 应保持 None，由渲染层 default_alignment="LEFT" 兜底，
        # 避免把左对齐的 title 错误渲染为居中。

    # 类型判断与分发
    # 优先用 has_table / has_text_frame 等方法判断，更可靠
    if shape.has_table:
        model.shape_type = "table"
        from ppt_transfor.parser.table import parse_table
        model.table = parse_table(shape, prs)
        return model

    # shape_type 枚举判断
    try:
        st = shape.shape_type
        st_name = st.name if hasattr(st, "name") else str(st)
    except Exception:
        st = None
        st_name = "UNKNOWN"

    if st == MSO_SHAPE_TYPE.GROUP:
        model.shape_type = "group"
        from ppt_transfor.parser.group import parse_group
        group_fields = parse_group(shape, slide, prs)
        model.children = group_fields["children"]
        if group_fields.get("child_offset") is not None:
            model.child_offset = group_fields["child_offset"]
        if group_fields.get("child_extent") is not None:
            model.child_extent = group_fields["child_extent"]
        return model

    if st == MSO_SHAPE_TYPE.PICTURE:
        model.shape_type = "picture"
        from ppt_transfor.parser.image import parse_picture
        pic_fields = parse_picture(shape)
        for k, v in pic_fields.items():
            setattr(model, k, v)
        # 图片也可能有文本（虽然少见）
        if shape.has_text_frame:
            from ppt_transfor.parser.text import parse_text_frame
            model.text = parse_text_frame(shape.text_frame, prs, inherited_props)
        return model

    if st == MSO_SHAPE_TYPE.CHART:
        model.shape_type = "chart"
        try:
            model.chart_xml = etree.tostring(shape._element, encoding="unicode")
        except Exception:
            pass
        try:
            model.chart_part = _parse_chart_part(shape, slide)
        except Exception:
            pass
        return model

    if st in (MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM) and not shape.has_text_frame:
        # 连接线（无文本框的线）
        model.shape_type = "connector"
        from ppt_transfor.parser.connector import parse_connector
        conn_fields = parse_connector(shape)
        for k, v in conn_fields.items():
            setattr(model, k, v)
        return model

    if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
        from ppt_transfor.parser.autoshape import parse_autoshape
        auto_fields = parse_autoshape(shape)

        # 无 auto_shape_type 且 line 有箭头 → 本质是带箭头的直线，重分类为 connector
        # 避免降级为文本框导致线条和箭头完全丢失
        if auto_fields.get("auto_shape_type") is None:
            line_model = getattr(model, "line", None)
            if line_model is not None and (
                line_model.head_arrow_type or line_model.tail_arrow_type
            ):
                model.shape_type = "connector"
                from ppt_transfor.parser.connector import parse_connector
                conn_fields = parse_connector(shape)
                for k, v in conn_fields.items():
                    setattr(model, k, v)
                return model

        model.shape_type = "auto_shape"
        for k, v in auto_fields.items():
            setattr(model, k, v)
        # 自选图形通常有文本
        if shape.has_text_frame:
            from ppt_transfor.parser.text import parse_text_frame
            model.text = parse_text_frame(shape.text_frame, prs, inherited_props)
        return model

    if st == MSO_SHAPE_TYPE.TEXT_BOX:
        model.shape_type = "text_box"
        if shape.has_text_frame:
            from ppt_transfor.parser.text import parse_text_frame
            model.text = parse_text_frame(shape.text_frame, prs, inherited_props)
        return model

    if st == MSO_SHAPE_TYPE.PLACEHOLDER:
        model.shape_type = "placeholder"
        if shape.has_text_frame:
            from ppt_transfor.parser.text import parse_text_frame
            model.text = parse_text_frame(shape.text_frame, prs, inherited_props)
        return model

    # 兜底：未知类型，尝试提取文本
    model.shape_type = st_name.lower() if st_name else "unknown"
    if shape.has_text_frame:
        from ppt_transfor.parser.text import parse_text_frame
        model.text = parse_text_frame(shape.text_frame, prs, inherited_props)

    return model
