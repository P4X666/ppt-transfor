"""幻灯片解析器：Slide → Slide 模型。

解析背景（含继承链）+ 遍历 shapes。
"""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE_TYPE

from ppt_transfor.models.schema import Color, Shape, Slide
from ppt_transfor.parser.shape import parse_shape
from ppt_transfor.utils.inheritance import resolve_background


def _rgb_to_luminance(hex_rgb: str) -> float:
    """RRGGBB → 相对亮度（0-1）。"""
    hex_rgb = hex_rgb.lstrip("#")
    if len(hex_rgb) != 6:
        return 0.5
    r = int(hex_rgb[0:2], 16) / 255
    g = int(hex_rgb[2:4], 16) / 255
    b = int(hex_rgb[4:6], 16) / 255
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _contrast_color(hex_bg: str) -> str:
    """根据背景色返回高对比度文本色（黑或白）。"""
    lum = _rgb_to_luminance(hex_bg)
    return "000000" if lum > 0.5 else "FFFFFF"


def _color_to_rgb(c: Color | None) -> str | None:
    """把 Color 模型转为 RRGGBB，主题/ scheme 返回 None。"""
    if c is None:
        return None
    if c.type == "rgb":
        val = c.value.lstrip("#")
        if len(val) == 6:
            return val
    return None


def _effective_bg_rgb(slide_model: Slide, shape_model) -> str | None:
    """从 Slide 模型和 Shape 模型计算有效背景 RGB。"""
    # 形状填充
    fill = getattr(shape_model, "fill", None)
    if fill is not None:
        if fill.type == "solid":
            rgb = _color_to_rgb(getattr(fill, "color", None))
            if rgb:
                return rgb
        # background 类型也跟随 slide bg，下方统一处理
    # 幻灯片背景
    bg = slide_model.background
    if bg is not None and bg.type == "solid":
        rgb = _color_to_rgb(bg.color)
        if rgb:
            return rgb
    return None


def _ensure_shape_text_visibility(
    slide_model: Slide,
    shape_model: Shape,
    threshold: float = 0.2,
) -> None:
    """修正单个 shape 内文本颜色，确保与背景有足够对比度。

    对透明填充（fill=None 或 type='none'）的 shape 不调整文本颜色：
    其实际背景可能是幻灯片背景，也可能是后层形状，模型层无法可靠判断，
    盲目调整可能把黑字改为白字后盖在浅色形状上，反而不可见。
    仅对显式纯色/背景填充的形状进行文本色修正。
    """
    fill = getattr(shape_model, "fill", None)
    fill_type = fill.type if fill is not None else None

    # 透明填充：不基于 slide bg 调整颜色，避免破坏后层浅色形状上的黑字
    if fill_type in (None, "none"):
        # 仍递归处理子形状，因为子形状可能有独立填充
        for child in getattr(shape_model, "children", []) or []:
            _ensure_shape_text_visibility(slide_model, child, threshold)
        return

    bg_hex = _effective_bg_rgb(slide_model, shape_model)
    if bg_hex is None:
        return
    bg_lum = _rgb_to_luminance(bg_hex)

    text_model = getattr(shape_model, "text", None)
    if text_model is not None:
        for para in text_model.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                text_rgb = _color_to_rgb(run.font.color)
                if text_rgb is None:
                    # 未显式设置颜色时，按默认黑色处理；若背景很深则设为白字
                    if bg_lum < 0.3:
                        run.font.color = Color(type="rgb", value="FFFFFF")
                    continue
                text_lum = _rgb_to_luminance(text_rgb)
                if abs(text_lum - bg_lum) < threshold:
                    new_hex = _contrast_color(bg_hex)
                    run.font.color = Color(type="rgb", value=new_hex)

    # 递归处理组合子形状
    for child in getattr(shape_model, "children", []) or []:
        _ensure_shape_text_visibility(slide_model, child, threshold)

    # 处理表格单元格
    table = getattr(shape_model, "table", None)
    if table is not None:
        for row in table.cells:
            for cell in row:
                cell_bg_hex = None
                if cell.fill is not None and cell.fill.type == "solid":
                    cell_bg_hex = _color_to_rgb(cell.fill.color)
                if cell_bg_hex is None:
                    cell_bg_hex = bg_hex
                cell_bg_lum = _rgb_to_luminance(cell_bg_hex)
                for para in cell.text.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        text_rgb = _color_to_rgb(run.font.color)
                        if text_rgb is None:
                            if cell_bg_lum < 0.3:
                                run.font.color = Color(type="rgb", value="FFFFFF")
                            continue
                        text_lum = _rgb_to_luminance(text_rgb)
                        if abs(text_lum - cell_bg_lum) < threshold:
                            new_hex = _contrast_color(cell_bg_hex)
                            run.font.color = Color(type="rgb", value=new_hex)


def _ensure_text_visibility(slide_model: Slide) -> None:
    """解析后处理：将低对比度文本调整为可见颜色。

    原始 PPT 中某些文本可能依赖图片/渐变背景或主题默认色，
    转换后若背景变为纯色且文本色与之接近，会导致黑底黑字/白底白字。
    这里在模型层统一修正，保证输出文件文字可见，同时让往返测试一致。
    本函数递归处理组合形状与表格单元格。
    """
    CONTRAST_THRESHOLD = 0.2
    for shape_model in slide_model.shapes:
        _ensure_shape_text_visibility(slide_model, shape_model, CONTRAST_THRESHOLD)


def _shape_dedup_key(model: Shape) -> str:
    """生成形状去重键。

    图片用完整 base64 的 MD5 去重（不同图片 PNG header 相同但数据不同，
    不能用前 N 字符判断）；非图片形状用类型+位置+填充组合键。
    """
    if model.shape_type == "picture" and model.data_base64:
        import hashlib

        return f"pic:{hashlib.md5(model.data_base64.encode()).hexdigest()}"
    # 非图片形状：用类型+位置+填充组合去重
    fill_str = ""
    if model.fill and model.fill.color and model.fill.color.value:
        fill_str = model.fill.color.value
    return (
        f"{model.shape_type}:"
        f"{model.left},{model.top},{model.width},{model.height}:"
        f"{fill_str}"
    )


def _collect_non_placeholder_shapes(container, prs=None) -> list:
    """收集 container 中所有非 placeholder 形状，返回原始 shape 对象列表。"""
    result = []
    try:
        for shape in container.shapes:
            try:
                st = shape.shape_type
                if st == MSO_SHAPE_TYPE.PLACEHOLDER:
                    continue
            except Exception:
                continue
            result.append(shape)
    except Exception:
        pass
    return result


def _has_fullscreen_picture(slide, prs=None) -> bool:
    """检查 slide 自己是否有覆盖整个幻灯片区域的图片。

    若 slide 有全屏 Image，PowerPoint 渲染时 slide 层的 Image 会覆盖 master 层，
    master 的装饰 shapes 不可见，无需合并 master shapes。
    典型场景：card.pptx slide[0] 有全屏透明 Image，master Rectangle 不应合并，
    否则透过透明 Image 看到 master Rectangle（灰色），而非 slide 背景（黑色）。
    """
    if prs is None:
        return False
    try:
        slide_width = prs.slide_width
        slide_height = prs.slide_height
    except Exception:
        return False

    for shape in slide.shapes:
        try:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            left = shape.left or 0
            top = shape.top or 0
            width = shape.width or 0
            height = shape.height or 0
            # 全屏条件：左上角接近原点，宽高覆盖 95%/90% 以上
            if left <= 0 and top <= slide_height * 0.1:
                if width >= slide_width * 0.95 and height >= slide_height * 0.90:
                    return True
        except Exception:
            continue
    return False


def _parse_inherited_pictures(slide, prs=None) -> list[Shape]:
    """解析布局/母版上继承的非 placeholder 形状，返回 Shape 列表。

    转换后 PPT 使用 Blank 布局，布局/母版继承的形状会丢失。
    此函数将非 placeholder 形状（图片/自选图形等）提取为幻灯片级形状，
    保证视觉保真。

    合并策略（模拟 OOXML 渲染层级 master < layout < slide）：
    - 若 layout 有非 placeholder shapes，说明该 layout 已定义自己的装饰
      （如全屏 Image 覆盖了 master 的灰色 Rectangle），只合并 layout 的 shapes。
    - 若 layout 无非 placeholder shapes 但 slide 自己有全屏 Image，
      slide 层的 Image 会覆盖 master 层，master 装饰不可见，不合并 master。
    - 若 layout 无非 placeholder shapes 且 slide 无全屏 Image，
      说明该 layout 依赖 master 的装饰，需合并 master 的 shapes（如左侧灰色 Rectangle）。
    """
    from ppt_transfor.parser.shape import parse_shape

    inherited: list[Shape] = []
    seen_keys: set[str] = set()

    try:
        layout = slide.slide_layout
    except Exception:
        return inherited

    # 优先收集 layout 的非 placeholder shapes
    layout_shapes = _collect_non_placeholder_shapes(layout, prs)

    if layout_shapes:
        # layout 有装饰 shapes → 只合并 layout（layout 已覆盖 master 的装饰）
        containers = [layout]
    elif _has_fullscreen_picture(slide, prs):
        # slide 有全屏 Image → slide 层覆盖 master 层，master 装饰不可见，不合并
        containers = []
    else:
        # layout 无装饰 shapes 且 slide 无全屏 Image → 合并 master（layout 依赖 master 的装饰）
        try:
            master = layout.slide_master
            containers = [master]
        except Exception:
            containers = []

    for container in containers:
        for shape in _collect_non_placeholder_shapes(container, prs):
            model = parse_shape(shape, None, prs)
            key = _shape_dedup_key(model)
            if key in seen_keys:
                continue
            seen_keys.add(key)
            inherited.append(model)

    return inherited


def parse_slide(slide, index: int, prs=None) -> Slide:
    """解析单页幻灯片

    Args:
        slide: python-pptx Slide 对象
        index: 幻灯片索引
        prs: 所属 Presentation 对象（传递给 shape 解析用于继承解析）
    """
    model = Slide(index=index)

    # 布局名
    try:
        model.layout_name = slide.slide_layout.name
    except Exception:
        model.layout_name = "Blank"

    # 背景：沿继承链解析（slide → layout → master），传入 prs 以固化主题色
    bg = resolve_background(slide, prs)
    if bg is not None:
        model.background = bg

    # 遍历形状
    for shape in slide.shapes:
        model.shapes.append(parse_shape(shape, slide, prs))

    # 前置布局/母版继承的图片（渲染在底层，保证视觉保真）
    # 转换后 PPT 使用 Blank 布局，不提取则布局/母版图片全部丢失
    inherited_pictures = _parse_inherited_pictures(slide, prs)
    if inherited_pictures:
        model.shapes = inherited_pictures + model.shapes

    # 后处理：修正低对比度文本，保证转换后可见
    _ensure_text_visibility(model)

    return model
