"""表格解析器：Table 形状 → Table 模型。

遍历单元格，记录合并（span_x/span_y）与文本。
"""

from __future__ import annotations

from pptx.oxml.ns import qn

from ppt_transfor.models.schema import CellBorder, Color, TableCell, Table
from ppt_transfor.parser.text import parse_text_frame
from ppt_transfor.utils.inheritance import extract_txbody_default_props


def parse_table(shape, prs=None) -> Table:
    """解析表格形状

    Args:
        shape: python-pptx GraphicFrame（含表格）
        prs: 所属 Presentation（用于主题色固化）
    """
    tbl = shape.table
    rows = len(tbl.rows)
    cols = len(tbl.columns)

    model = Table(rows=rows, cols=cols)

    # 首行表头标记（通过样式推断，简化处理）
    model.first_row_header = bool(getattr(tbl, "first_row", False))

    # 解析 tblPr 属性和 tableStyleId，保留原始表格样式引用
    # 避免渲染时 add_table 使用默认蓝色样式
    try:
        tbl_el = shape._element
        tbl_pr = tbl_el.find(".//" + qn("a:tblPr"))
        if tbl_pr is not None:
            if tbl_pr.get("firstCol") == "1":
                model.first_col = True
            if tbl_pr.get("lastRow") == "1":
                model.last_row = True
            if tbl_pr.get("bandRow") == "1":
                model.band_row = True
            style_id_el = tbl_pr.find(qn("a:tableStyleId"))
            if style_id_el is not None and style_id_el.text:
                model.style_id = style_id_el.text.strip()
    except Exception:
        pass

    # 解析列宽（EMU 整数），保留原始列宽避免 add_table 默认均分
    try:
        col_widths = []
        for col in tbl.columns:
            col_widths.append(int(col.width))
        model.col_widths = col_widths
    except Exception:
        pass

    # 单元格矩阵
    cells = []
    for row_idx in range(rows):
        row_cells = []
        for col_idx in range(cols):
            cell = tbl.cell(row_idx, col_idx)
            cell_model = TableCell()

            # 文本：传入单元格默认样式，避免对齐/字号丢失
            cell_defaults = extract_txbody_default_props(cell.text_frame._element, prs)
            cell_model.text = parse_text_frame(cell.text_frame, prs, cell_defaults)

            # 合并信息：span_x 为横向合并数，span_y 为纵向合并数
            # python-pptx 通过 cell.span_x / cell.span_y 获取
            try:
                cell_model.span_x = int(cell.span_x)
            except Exception:
                cell_model.span_x = 1
            try:
                cell_model.span_y = int(cell.span_y)
            except Exception:
                cell_model.span_y = 1

            # 单元格填充：解析 <a:tcPr> 下的 solidFill 等，主题色固化为 RGB
            # 未解析填充会导致表头颜色丢失（add_table 默认蓝色样式）
            try:
                from ppt_transfor.parser.shape import _parse_fill

                cell_fill = _parse_fill(cell.fill, cell._tc, prs)
                if cell_fill is not None:
                    cell_model.fill = cell_fill
            except Exception:
                pass

            # 单元格边框：解析 <a:tcPr> 下的 lnL/lnR/lnT/lnB
            # 保留原始边框信息，避免 add_table 默认添加边框
            try:
                tc_pr = cell._tc.find(qn("a:tcPr"))
                if tc_pr is not None:
                    for side, tag in (("left", "lnL"), ("right", "lnR"),
                                      ("top", "lnT"), ("bottom", "lnB")):
                        ln_el = tc_pr.find(qn(f"a:{tag}"))
                        if ln_el is not None:
                            border = CellBorder()
                            w = ln_el.get("w")
                            if w:
                                try:
                                    border.width = int(w)
                                except (ValueError, TypeError):
                                    pass
                            # 检测 noFill
                            no_fill_el = ln_el.find(qn("a:noFill"))
                            if no_fill_el is not None:
                                border.no_fill = True
                            # 检测 solidFill
                            solid_fill = ln_el.find(qn("a:solidFill"))
                            if solid_fill is not None:
                                srgb = solid_fill.find(qn("a:srgbClr"))
                                if srgb is not None and srgb.get("val"):
                                    border.color = Color(type="rgb", value=srgb.get("val"))
                            # 有 width 但无 solidFill 且无 noFill → 标记为 no_fill
                            if border.width is not None and border.color is None and not border.no_fill:
                                border.no_fill = True
                            setattr(cell_model, f"border_{side}", border)
            except Exception:
                pass

            row_cells.append(cell_model)
        cells.append(row_cells)

    model.cells = cells
    return model
