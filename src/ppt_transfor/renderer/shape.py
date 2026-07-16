"""形状渲染器：通用属性回写 + 按类型分发。

通用属性：位置/旋转/填充/边框/阴影
类型分发：table/group/picture/connector/autoshape/text_box/placeholder
"""

from __future__ import annotations

from pptx.oxml.ns import qn
from pptx.util import Emu

from ppt_transfor.models.schema import Color, Fill, Line, Shape
from ppt_transfor.utils.color import apply_color


def _set_no_fill(fill) -> None:
    """显式设置形状填充为 <a:noFill/>。

    python-pptx 未暴露 fill.none() API，需直接操作 spPr 下的 XML。
    """
    from lxml import etree

    spPr = fill._xPr
    # 移除现有填充子元素
    for tag in ("noFill", "solidFill", "gradFill", "blipFill", "pattFill", "grpFill"):
        child = spPr.find(qn(f"a:{tag}"))
        if child is not None:
            spPr.remove(child)
    # 插入 noFill（追加即可，渲染器按标签查找，不严格依赖顺序）
    no_fill = etree.Element(qn("a:noFill"))
    spPr.append(no_fill)


def _apply_gradient_fill(fill, fill_model: Fill) -> None:
    """将 Fill 模型中的渐变 stops / 类型 / 角度回写到 XML。

    python-pptx 未暴露 gradient stops API，需直接操作 <a:gradFill>。
    """
    from lxml import etree

    if not fill_model.gradient_stops:
        return

    spPr = fill._xPr
    grad_fill = spPr.find(qn("a:gradFill"))
    if grad_fill is None:
        return

    # 清除现有 stops 与方向定义
    for tag in ("gsLst", "lin", "path", "rect", "tileRect"):
        child = grad_fill.find(qn(f"a:{tag}"))
        if child is not None:
            grad_fill.remove(child)

    # stops：按 OpenXML 标准，颜色直接作为 <a:gs> 的子元素，不包 <a:solidFill>
    gs_lst = etree.Element(qn("a:gsLst"))
    for stop in fill_model.gradient_stops:
        gs = etree.SubElement(gs_lst, qn("a:gs"))
        # pos 为 0~100000 的整数
        gs.set("pos", str(int(round(stop.position * 100000))))
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        val = stop.color.value.lstrip("#")
        if len(val) != 6:
            val = "FFFFFF"
        srgb.set("val", val.upper())
    grad_fill.append(gs_lst)

    # 方向定义
    gtype = fill_model.gradient_type
    if gtype == "linear":
        lin = etree.Element(qn("a:lin"))
        angle = fill_model.gradient_angle if fill_model.gradient_angle is not None else 5400000
        lin.set("ang", str(int(angle)))
        lin.set("scaled", "1")
        grad_fill.append(lin)
    elif gtype in ("circle", "shape"):
        # path 渐变：path 属性常见值 circle / shape
        path = etree.Element(qn("a:path"))
        path.set("path", gtype)
        grad_fill.append(path)
    elif gtype == "rect":
        rect = etree.Element(qn("a:rect"))
        # 默认值与 OpenXML 一致
        rect.set("l", "50000")
        rect.set("t", "50000")
        rect.set("r", "50000")
        rect.set("b", "50000")
        grad_fill.append(rect)


def _apply_fill(fill, fill_model: Fill | None, slide_bg_color: Color | None = None) -> None:
    """回写填充

    Args:
        fill: python-pptx FillFormat
        fill_model: Fill 模型
        slide_bg_color: 当前幻灯片背景色（solid 时），用于 BACKGROUND 填充兜底
    """
    if fill_model is None:
        return

    if fill_model.type == "none":
        # 显式透明：在 spPr 下写入 <a:noFill/>
        try:
            _set_no_fill(fill)
        except Exception:
            pass
        return

    if fill_model.type == "background":
        # 使用幻灯片背景填充，保持与原始文件一致的 fill 类型
        try:
            fill.background()
        except Exception:
            pass
        return

    if fill_model.type == "solid":
        # 无颜色时不调用 solid()，防止默认黑色
        if fill_model.color is None:
            return
        try:
            fill.solid()
            apply_color(fill.fore_color, fill_model.color)
        except Exception:
            pass
        return

    if fill_model.type == "gradient":
        # 渐变：先调用 python-pptx API 创建 gradFill 骨架，再用 XML 写入 stops 和方向
        try:
            fill.gradient()
            _apply_gradient_fill(fill, fill_model)
        except Exception:
            pass
        return

    # pattern/picture 等高级填充暂不支持回写
    # 后续可通过 xml_helper 直接操作 XML 扩展


def _apply_line(line, line_model: Line | None) -> None:
    """回写边框"""
    if line_model is None:
        return

    if line_model.width is not None:
        try:
            line.width = Emu(line_model.width)
        except Exception:
            pass

    # 无填充线条：显式设置 <a:noFill/>，避免继承 p:style 或主题默认色（蓝色）
    # 必须在 width 设置之后执行，确保 <a:ln> 已存在
    if line_model.no_fill:
        try:
            from lxml import etree

            ln_el = line._ln
            if ln_el is not None:
                # 移除现有 solidFill
                existing = ln_el.find(qn("a:solidFill"))
                if existing is not None:
                    ln_el.remove(existing)
                # noFill 插入到第一个位置，确保 schema 顺序正确
                # OpenXML 规范要求 <a:ln> 子元素顺序：noFill/solidFill → dash → round/bevel/miter → headEnd → tailEnd
                no_fill_el = ln_el.find(qn("a:noFill"))
                if no_fill_el is not None:
                    ln_el.remove(no_fill_el)
                no_fill_el = etree.Element(qn("a:noFill"))
                ln_el.insert(0, no_fill_el)
        except Exception:
            pass
    elif line_model.color is not None:
        try:
            apply_color(line.color, line_model.color)
        except Exception:
            pass

    if line_model.dash is not None:
        try:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            line.dash_style = MSO_LINE_DASH_STYLE[line_model.dash]
        except (KeyError, ValueError, Exception):
            pass

    # 箭头线端点：回写 <a:headEnd>/<a:tailEnd> 到 <a:ln>
    # 必须在 width/color/dash 设置之后执行，确保 <a:ln> 已存在
    if line_model.head_arrow_type or line_model.tail_arrow_type:
        try:
            from lxml import etree

            ln_el = line._ln
            if ln_el is not None:
                # 清除现有箭头
                for tag in ("headEnd", "tailEnd"):
                    existing = ln_el.find(qn(f"a:{tag}"))
                    if existing is not None:
                        ln_el.remove(existing)
                # 写入新箭头
                if line_model.head_arrow_type:
                    head = etree.SubElement(ln_el, qn("a:headEnd"))
                    head.set("type", line_model.head_arrow_type)
                if line_model.tail_arrow_type:
                    tail = etree.SubElement(ln_el, qn("a:tailEnd"))
                    tail.set("type", line_model.tail_arrow_type)
        except Exception:
            pass


def _apply_common_props(shape, model: Shape, slide_bg_color: Color | None = None) -> None:
    """回写通用属性（位置/旋转/填充/边框/阴影）"""
    # 位置与尺寸
    if model.left is not None:
        try:
            shape.left = Emu(model.left)
        except Exception:
            pass
    if model.top is not None:
        try:
            shape.top = Emu(model.top)
        except Exception:
            pass
    if model.width is not None:
        try:
            shape.width = Emu(model.width)
        except Exception:
            pass
    if model.height is not None:
        try:
            shape.height = Emu(model.height)
        except Exception:
            pass

    # 旋转
    if model.rotation:
        try:
            shape.rotation = model.rotation
        except Exception:
            pass

    # 填充
    try:
        _apply_fill(shape.fill, model.fill, slide_bg_color)
    except Exception:
        pass

    # 边框
    try:
        _apply_line(shape.line, model.line)
    except Exception:
        pass

    # 阴影：若模型未指定，显式关闭继承，避免 add_shape 默认继承主题阴影
    if model.shadow is not None:
        try:
            shape.shadow.inherit = model.shadow.enabled
        except Exception:
            pass
    else:
        try:
            shape.shadow.inherit = False
        except Exception:
            pass

    # 移除 add_shape/add_connector 自动创建的 <p:style> 元素（原始 PPT 没有）
    # p:style 的 lnRef 会引用主题线条样式，可能导致渲染器显示主题边框
    # 即使 spPr 的 <a:ln> 有 noFill，某些渲染器仍可能因 p:style 存在而回退到主题边框
    # 注意：add_textbox 不创建 p:style，所以 text_box 不受影响
    if not getattr(model, "has_style", False):
        try:
            p_style = shape._element.find(qn("p:style"))
            if p_style is not None:
                shape._element.remove(p_style)
        except Exception:
            pass
    elif getattr(model, "style_xml", None):
        # 原始 PPT 有 p:style，但 add_textbox 不创建 p:style，需要从原始 PPT 复制
        try:
            from lxml import etree
            # 移除现有的 p:style（如果有）
            existing = shape._element.find(qn("p:style"))
            if existing is not None:
                shape._element.remove(existing)
            # 解析并插入 p:style
            new_style = etree.fromstring(model.style_xml)
            # OOXML schema 要求 sp 子元素顺序: nvSpPr → spPr → style → txBody
            # style 必须在 spPr 之后（addprevious 会导致 style 在 spPr 之前，违反 schema）
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is not None:
                spPr.addnext(new_style)
            else:
                shape._element.append(new_style)
        except Exception:
            pass


def render_shape(container, model: Shape, slide_bg_color: Color | None = None):
    """渲染单个形状到容器（slide 或 group），按类型分发。

    Args:
        container: python-pptx Slide 或 GroupShape
        model: Shape 模型
        slide_bg_color: 当前幻灯片背景色（solid 时），用于 BACKGROUND 填充兜底

    Returns:
        渲染后的 python-pptx 形状对象
    """
    shape = None

    if model.shape_type == "table":
        from ppt_transfor.renderer.table import render_table
        shape = render_table(container, model)
    elif model.shape_type == "group":
        from ppt_transfor.renderer.group import render_group
        shape = render_group(container, model, slide_bg_color)
    elif model.shape_type == "picture":
        from ppt_transfor.renderer.image import render_picture
        shape = render_picture(container, model)
    elif model.shape_type == "connector":
        from ppt_transfor.renderer.connector import render_connector
        shape = render_connector(container, model)
    elif model.shape_type == "auto_shape":
        from ppt_transfor.renderer.autoshape import render_autoshape
        shape = render_autoshape(container, model)
        # auto_shape_type 为 None 或无法识别时，降级为文本框保持位置和尺寸
        # 文本由后续统一逻辑渲染，避免重复
        if shape is None:
            shape = container.shapes.add_textbox(
                Emu(model.left) if model.left is not None else None,
                Emu(model.top) if model.top is not None else None,
                Emu(model.width) if model.width is not None else None,
                Emu(model.height) if model.height is not None else None,
            )
    elif model.shape_type in ("text_box", "placeholder"):
        # 文本框：add_textbox 后回写文本，默认左对齐、默认自动换行
        from ppt_transfor.renderer.text import render_text_frame
        shape = container.shapes.add_textbox(
            Emu(model.left) if model.left is not None else None,
            Emu(model.top) if model.top is not None else None,
            Emu(model.width) if model.width is not None else None,
            Emu(model.height) if model.height is not None else None,
        )
        if model.text is not None:
            render_text_frame(
                shape.text_frame,
                model.text,
                default_alignment="LEFT",
                default_word_wrap=True,
            )
    elif model.shape_type == "chart":
        # chart 通过保存后的 zip 后处理插入，这里不生成占位文本框
        return None
    else:
        # 兜底：未知类型，渲染占位文本框保持位置和尺寸
        from ppt_transfor.renderer.text import render_text_frame
        shape = container.shapes.add_textbox(
            Emu(model.left) if model.left is not None else None,
            Emu(model.top) if model.top is not None else None,
            Emu(model.width) if model.width is not None else None,
            Emu(model.height) if model.height is not None else None,
        )
        if model.text is not None:
            render_text_frame(
                shape.text_frame,
                model.text,
                default_alignment="LEFT",
                default_word_wrap=True,
            )

    # 回写通用属性（图片/表格的位置已在 add_xxx 时设置，这里再次确保）
    if shape is not None:
        _apply_common_props(shape, model, slide_bg_color)

        # 自选图形/文本框的文本回写（auto_shape 在 add_shape 时已创建空文本框）
        if model.shape_type == "auto_shape" and model.text is not None:
            from ppt_transfor.renderer.text import render_text_frame
            try:
                # auto_shape 的默认段落对齐为居中，模型未显式指定时兜底居中
                render_text_frame(
                    shape.text_frame,
                    model.text,
                    default_alignment="CENTER",
                    default_word_wrap=True,
                )
            except Exception:
                pass

    return shape
