"""表格渲染器：Table 模型 → Table 形状。

add_table 后填充单元格，处理合并。
"""

from __future__ import annotations

from lxml import etree
from pptx.oxml.ns import qn
from pptx.util import Emu

from ppt_transfor.models.schema import Shape
from ppt_transfor.renderer.text import render_text_frame

# OpenXML schema 中 <a:tcPr> 子元素的顺序
# lnL, lnR, lnT, lnB, lnTlToBr, lnBlToTr, cell3D, fill(noFill/solidFill/...), headers, extLst
_TC_PR_CHILD_ORDER = [
    "lnL", "lnR", "lnT", "lnB", "lnTlToBr", "lnBlToTr", "cell3D",
    "noFill", "solidFill", "gradFill", "blipFill", "pattFill", "grpFill",
    "headers", "extLst",
]


def _insert_tc_pr_child(tc_pr, element, tag_name: str) -> None:
    """按 OpenXML schema 顺序插入 tcPr 子元素。

    确保边框元素（lnL/lnR/lnT/lnB）在填充元素（solidFill 等）之前，
    避免 PowerPoint 因 schema 顺序错误忽略边框定义。
    """
    insert_idx = _TC_PR_CHILD_ORDER.index(tag_name) if tag_name in _TC_PR_CHILD_ORDER else len(_TC_PR_CHILD_ORDER)

    for i, child in enumerate(tc_pr):
        child_name = etree.QName(child).localname
        if child_name in _TC_PR_CHILD_ORDER:
            child_idx = _TC_PR_CHILD_ORDER.index(child_name)
            if child_idx > insert_idx:
                tc_pr.insert(i, element)
                return
    tc_pr.append(element)


def render_table(slide, model: Shape):
    """渲染表格形状，返回 python-pptx GraphicFrame 对象。

    Args:
        slide: python-pptx Slide
        model: Shape 模型（shape_type == "table"）

    Returns:
        GraphicFrame 对象
    """
    tbl_model = model.table
    if tbl_model is None:
        return None

    rows = tbl_model.rows
    cols = tbl_model.cols

    graphic_frame = slide.shapes.add_table(
        rows,
        cols,
        Emu(model.left) if model.left is not None else None,
        Emu(model.top) if model.top is not None else None,
        Emu(model.width) if model.width is not None else None,
        Emu(model.height) if model.height is not None else None,
    )
    tbl = graphic_frame.table

    # 首行表头：显式设置，避免 add_table 默认 True
    try:
        tbl.first_row = tbl_model.first_row_header
    except Exception:
        pass

    # 回写 tblPr 属性和 tableStyleId，保留原始表格样式引用
    # 避免使用 add_table 默认蓝色样式和默认边框
    try:
        tbl_el = graphic_frame._element
        tbl_pr = tbl_el.find(".//" + qn("a:tblPr"))
        if tbl_pr is not None:
            # 显式设置 firstCol/lastRow/bandRow，覆盖 add_table 默认值
            tbl_pr.set("firstCol", "1" if tbl_model.first_col else "0")
            tbl_pr.set("lastRow", "1" if tbl_model.last_row else "0")
            tbl_pr.set("bandRow", "1" if tbl_model.band_row else "0")
            # tableStyleId
            if tbl_model.style_id:
                existing = tbl_pr.find(qn("a:tableStyleId"))
                if existing is not None:
                    tbl_pr.remove(existing)
                style_id_el = etree.SubElement(tbl_pr, qn("a:tableStyleId"))
                style_id_el.text = tbl_model.style_id
    except Exception:
        pass

    # 回写列宽，保留原始列宽避免 add_table 默认均分
    if tbl_model.col_widths:
        try:
            for idx, width in enumerate(tbl_model.col_widths):
                if idx < len(tbl.columns):
                    tbl.columns[idx].width = Emu(width)
        except Exception:
            pass

    # 填充单元格
    for row_idx in range(rows):
        for col_idx in range(cols):
            if row_idx >= len(tbl_model.cells) or col_idx >= len(tbl_model.cells[row_idx]):
                continue
            cell_model = tbl_model.cells[row_idx][col_idx]
            cell = tbl.cell(row_idx, col_idx)

            # 合并：span_x/span_y > 1 时需要 merge
            # 注意：python-pptx 的 merge 会将目标单元格合并到当前单元格
            if cell_model.span_x > 1 or cell_model.span_y > 1:
                try:
                    end_row = row_idx + cell_model.span_y - 1
                    end_col = col_idx + cell_model.span_x - 1
                    if end_row < rows and end_col < cols:
                        cell.merge(tbl.cell(end_row, end_col))
                except Exception:
                    pass

            # 文本
            render_text_frame(cell.text_frame, cell_model.text)

            # 单元格填充：回写到 <a:tcPr>，覆盖 add_table 默认蓝色样式
            if cell_model.fill is not None:
                try:
                    from ppt_transfor.renderer.shape import _apply_fill

                    _apply_fill(cell.fill, cell_model.fill, None)
                except Exception:
                    pass

            # 单元格边框：回写 lnL/lnR/lnT/lnB 到 <a:tcPr>
            # 保留原始边框信息，避免 add_table 默认添加边框
            # 按 OpenXML schema 顺序插入：lnL, lnR, lnT, lnB 必须在 fill 之前
            try:
                tc_pr = cell._tc.find(qn("a:tcPr"))
                if tc_pr is None:
                    tc_pr = etree.SubElement(cell._tc, qn("a:tcPr"))
                for side, tag in (("left", "lnL"), ("right", "lnR"),
                                  ("top", "lnT"), ("bottom", "lnB")):
                    border = getattr(cell_model, f"border_{side}", None)
                    # 移除现有边框
                    existing = tc_pr.find(qn(f"a:{tag}"))
                    if existing is not None:
                        tc_pr.remove(existing)
                    if border is not None:
                        ln_el = etree.Element(qn(f"a:{tag}"))
                        if border.width is not None:
                            ln_el.set("w", str(border.width))
                        if border.no_fill:
                            # noFill 插入到 ln 的第一个位置，确保 schema 顺序
                            etree.SubElement(ln_el, qn("a:noFill"))
                            # noFill 应在 solidFill 之前，但 SubElement 追加到末尾
                            # 由于 ln_el 是新建的，noFill 是第一个子元素，顺序正确
                        elif border.color is not None:
                            solid = etree.SubElement(ln_el, qn("a:solidFill"))
                            srgb = etree.SubElement(solid, qn("a:srgbClr"))
                            srgb.set("val", border.color.value.lstrip("#").upper())
                        # 按 schema 顺序插入 tcPr（lnL/lnR/lnT/lnB 在 fill 之前）
                        _insert_tc_pr_child(tc_pr, ln_el, tag)
            except Exception:
                pass

    return graphic_frame
