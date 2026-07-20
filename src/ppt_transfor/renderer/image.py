"""图片渲染器：Picture 模型 → Picture 形状。

优先从 image_path 读取文件，降级从 data_base64 读取（向后兼容旧 JSON）。
输出 PPTX 仍嵌入二进制（add_picture），保持自包含可分发。
"""

from __future__ import annotations

import base64
import io
from pathlib import Path

from pptx.util import Emu

from ppt_transfor.models.schema import Shape


def render_picture(slide, model: Shape, base_dir: Path | None = None):
    """渲染图片形状，返回 python-pptx Picture 对象。

    Args:
        slide: python-pptx Slide 或 Group 容器
        model: Shape 模型（shape_type == "picture"）
        base_dir: image_path 的基准目录（如 out/），None 时用 cwd

    Returns:
        Picture 对象；若图片数据缺失则返回 None
    """
    image_stream = _get_image_stream(model, base_dir)
    if image_stream is None:
        return None

    pic = slide.shapes.add_picture(
        image_stream,
        Emu(model.left) if model.left is not None else None,
        Emu(model.top) if model.top is not None else None,
        Emu(model.width) if model.width is not None else None,
        Emu(model.height) if model.height is not None else None,
    )

    # 裁剪
    if model.crop is not None:
        try:
            pic.crop_left = model.crop.left
            pic.crop_top = model.crop.top
            pic.crop_right = model.crop.right
            pic.crop_bottom = model.crop.bottom
        except Exception:
            pass

    return pic


def _get_image_stream(model: Shape, base_dir: Path | None = None) -> io.BytesIO | None:
    """从 image_path 或 data_base64 获取图片字节流。

    优先 image_path（新格式，读取文件），降级 data_base64（旧格式，解码 base64）。
    两者都失败返回 None。
    """
    # 优先从文件读取（新格式）
    if model.image_path:
        try:
            base = Path(base_dir) if base_dir else Path.cwd()
            file_path = base / model.image_path
            if file_path.exists():
                blob = file_path.read_bytes()
                return io.BytesIO(blob)
        except Exception:
            pass

    # 降级从 base64 读取（向后兼容旧 JSON）
    if model.data_base64:
        try:
            blob = base64.b64decode(model.data_base64)
            return io.BytesIO(blob)
        except Exception:
            pass

    return None
