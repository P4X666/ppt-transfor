from pptx import Presentation
prs = Presentation('input/conent_page_component_card.pptx')
for si, slide in enumerate(prs.slides):
    print(f'Slide {si}')
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        tf = shape.text_frame
        print(f"  Shape {shape.name}: type={shape.shape_type} w={shape.width} h={shape.height}")
        print(f"    word_wrap={tf.word_wrap} auto_size={tf.auto_size} vanchor={tf.vertical_anchor}")
        for pi, para in enumerate(tf.paragraphs):
            print(f"    Para[{pi}]: align={para.alignment} level={para.level}")
            for ri, run in enumerate(para.runs):
                ct = run.font.color.type
                ctn = ct.name if hasattr(ct, 'name') else str(ct) if ct else None
                print(f"      Run[{ri}]: text={run.text[:30]!r} size={run.font.size} color={ctn}")
