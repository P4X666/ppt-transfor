from lxml import etree
from pptx import Presentation
prs = Presentation('input/content_page_component_calendar.pptx')
slide = prs.slides[0]

for s in slide.shapes:
    if s.name == 'Group 7':
        for child in s.shapes:
            if child.name == 'Rectangle 80':
                print(etree.tostring(child._element, pretty_print=True, encoding='unicode')[:2000])
                break
        break
