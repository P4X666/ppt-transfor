"""检测转换后 PPT 中的文本可见性问题（如黑字黑底）。

对 input 下所有 pptx 做转换，然后分析每个 run 的字体颜色与所在形状/幻灯片背景
的对比度，低对比度则报告。
"""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent / "src"))

from pptx import Presentation as PptxPresentation
from ppt_transfor.parser.presentation import parse_presentation
from ppt_transfor.renderer.presentation import render_presentation


def rgb_to_luminance(hex_rgb: str) -> float:
    """把 RRGGBB 转相对亮度（0-1），用于判断深浅。"""
    hex_rgb = hex_rgb.lstrip("#")
    if len(hex_rgb) != 6:
        return 0.5
    r = int(hex_rgb[0:2], 16) / 255
    g = int(hex_rgb[2:4], 16) / 255
    b = int(hex_rgb[4:6], 16) / 255
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def get_slide_bg_rgb(slide) -> str | None:
    """取幻灯片背景 RGB，仅处理 SOLID。"""
    try:
        fill = slide.background.fill
        ft = fill.type
        ft_name = ft.name if hasattr(ft, "name") else str(ft)
        if ft_name == "SOLID":
            rgb = fill.fore_color.rgb
            if rgb:
                return str(rgb)
    except Exception:
        pass
    return None


def get_shape_fill_rgb(shape) -> str | None:
    """取形状填充 RGB，BACKGROUND/None 返回 None 让上层用 slide bg。"""
    try:
        fill = shape.fill
        ft = fill.type
        ft_name = ft.name if hasattr(ft, "name") else str(ft)
        if ft_name == "SOLID":
            rgb = fill.fore_color.rgb
            if rgb:
                return str(rgb)
    except Exception:
        pass
    return None


def get_run_color_rgb(run) -> str | None:
    """取 run 字体颜色 RGB，None/SCHEME 视为未解析。"""
    try:
        c = run.font.color
        ct = c.type
        ct_name = ct.name if hasattr(ct, "name") else str(ct)
        if ct_name == "RGB":
            rgb = c.rgb
            if rgb:
                return str(rgb)
    except Exception:
        pass
    return None


def is_low_contrast(text_hex: str, bg_hex: str) -> bool:
    """文本与背景亮度差是否过低（< 0.2）。"""
    tl = rgb_to_luminance(text_hex)
    bl = rgb_to_luminance(bg_hex)
    return abs(tl - bl) < 0.2


def analyze(pptx_path: Path):
    prs = PptxPresentation(str(pptx_path))
    issues = []
    for si, slide in enumerate(prs.slides):
        slide_bg = get_slide_bg_rgb(slide) or "FFFFFF"
        for shape in slide.shapes:
            shape_bg = get_shape_fill_rgb(shape) or slide_bg
            if not shape.has_text_frame:
                continue
            for pi, para in enumerate(shape.text_frame.paragraphs):
                for ri, run in enumerate(para.runs):
                    text = run.text.strip()
                    if not text:
                        continue
                    color = get_run_color_rgb(run)
                    if color is None:
                        # 未解析到颜色，可能是继承默认（在黑底上通常为白）
                        continue
                    if is_low_contrast(color, shape_bg):
                        issues.append({
                            "slide": si,
                            "shape": shape.name,
                            "para": pi,
                            "run": ri,
                            "text": text[:30],
                            "text_color": color,
                            "bg_color": shape_bg,
                        })
    return issues


def main():
    input_dir = Path("input")
    out_dir = Path("out/pptx")
    out_dir.mkdir(parents=True, exist_ok=True)

    files = sorted(input_dir.glob("*.pptx"))
    files = [f for f in files if not f.name.startswith("~$")]

    all_issues = {}
    for pptx in files:
        converted = out_dir / pptx.name
        try:
            model = parse_presentation(str(pptx))
            render_presentation(model, str(converted))
        except Exception as e:
            print(f"转换失败 {pptx.name}: {e}")
            continue
        issues = analyze(converted)
        if issues:
            all_issues[pptx.name] = issues
            print(f"⚠️ {pptx.name}: 发现 {len(issues)} 处低对比度文本")
            for issue in issues[:5]:
                print(f"    slide={issue['slide']} shape={issue['shape']} text='{issue['text']}' text_color=#{issue['text_color']} bg=#{issue['bg_color']}")
            if len(issues) > 5:
                print(f"    ... 还有 {len(issues) - 5} 处")
        else:
            print(f"✅ {pptx.name}: 无低对比度文本")

    print(f"\n共 {len(all_issues)} 个文件存在潜在可见性问题")


if __name__ == "__main__":
    main()
