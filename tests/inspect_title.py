from pptx import Presentation
from lxml import etree

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

def show_title_ph(path):
    prs = Presentation(path)
    for si, slide in enumerate(prs.slides):
        if not slide.shapes:
            continue
        shape = slide.shapes[0]
        if shape.shape_type.name != "PLACEHOLDER":
            continue
        print(f"\nSlide {si}, shape {shape.name}, placeholder idx={shape.placeholder_format.idx}")
        # 段落对齐
        if shape.has_text_frame:
            tf = shape.text_frame
            for pi, para in enumerate(tf.paragraphs):
                print(f"  para[{pi}] alignment API={para.alignment}")
                pPr = para._element.find(f"{{{NS_A}}}pPr")
                if pPr is not None:
                    print(f"    pPr@algn={pPr.get('algn')}")
        # layout placeholder
        layout = slide.slide_layout
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == shape.placeholder_format.idx:
                print(f"  layout placeholder: {ph.name}")
                if ph.has_text_frame:
                    for pi, para in enumerate(ph.text_frame.paragraphs):
                        print(f"    para[{pi}] alignment API={para.alignment}")
                        pPr = para._element.find(f"{{{NS_A}}}pPr")
                        if pPr is not None:
                            print(f"      pPr@algn={pPr.get('algn')}")
                # lstStyle
                txBody = ph._element.find(f".//{{{NS_P}}}txBody")
                if txBody is not None:
                    lstStyle = txBody.find(f"{{{NS_A}}}lstStyle")
                    if lstStyle is not None:
                        for lvl in lstStyle:
                            tag = etree.QName(lvl).localname
                            algn = lvl.get("algn")
                            print(f"    lstStyle/{tag}@algn={algn}")
        # master placeholder
        master = layout.slide_master
        for ph in master.placeholders:
            if ph.placeholder_format.idx == shape.placeholder_format.idx:
                print(f"  master placeholder: {ph.name}")
                if ph.has_text_frame:
                    for pi, para in enumerate(ph.text_frame.paragraphs):
                        print(f"    para[{pi}] alignment API={para.alignment}")
                        pPr = para._element.find(f"{{{NS_A}}}pPr")
                        if pPr is not None:
                            print(f"      pPr@algn={pPr.get('algn')}")
                txBody = ph._element.find(f".//{{{NS_P}}}txBody")
                if txBody is not None:
                    lstStyle = txBody.find(f"{{{NS_A}}}lstStyle")
                    if lstStyle is not None:
                        for lvl in lstStyle:
                            tag = etree.QName(lvl).localname
                            algn = lvl.get("algn")
                            print(f"    lstStyle/{tag}@algn={algn}")

show_title_ph('input/conent_page_component_card.pptx')
show_title_ph('input/content_page_component_calendar.pptx')
