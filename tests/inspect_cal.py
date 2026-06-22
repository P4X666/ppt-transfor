from pptx import Presentation
prs = Presentation('input/content_page_component_calendar.pptx')

slide = prs.slides[0]
print('slide bg fill type:', slide.background.fill.type)
layout = slide.slide_layout
master = layout.slide_master
for name, bg in [('layout', layout.background), ('master', master.background)]:
    try:
        ft = bg.fill.type
        fn = ft.name if hasattr(ft, 'name') else str(ft)
    except Exception:
        fn = 'ERR'
    try:
        rgb = bg.fill.fore_color.rgb if bg.fill.fore_color else None
    except Exception:
        rgb = None
    print(f'{name} bg fill={fn} rgb={rgb}')

from lxml import etree

def show(elem, indent=0):
    if elem is None:
        return
    tag = etree.QName(elem).localname
    attrs = ' '.join(f'{k}="{v}"' for k, v in elem.attrib.items())
    print('  ' * indent + f'<{tag} {attrs}>'.rstrip())
    for child in elem:
        show(child, indent + 1)

print('\n--- find month labels in group ---')
def walk_shapes(shapes, indent=0):
    for shape in shapes:
        name = getattr(shape, 'name', '')
        if 'January' in name or 'February' in name or 'March' in name or 'Month' in name:
            print('  ' * indent + f'Found {name}:')
            try:
                ft = shape.fill.type
                fn = ft.name if hasattr(ft, 'name') else str(ft)
            except Exception:
                fn = 'ERR'
            try:
                rgb = shape.fill.fore_color.rgb if shape.fill.fore_color else None
            except Exception:
                rgb = None
            print(f"    fill={fn} rgb={rgb} left={shape.left} top={shape.top} w={shape.width} h={shape.height}")
            if hasattr(shape, '_element'):
                txbody = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}txBody')
                if txbody is not None:
                    show(txbody)
        if hasattr(shape, 'shapes'):
            walk_shapes(shape.shapes, indent + 1)

walk_shapes(slide.shapes)
