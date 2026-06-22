from lxml import etree
from pptx import Presentation

prs = Presentation('input/conent_page_component_card.pptx')
nsmap = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}

def show(elem, indent=0):
    if elem is None:
        return
    tag = etree.QName(elem).localname
    attrs = ' '.join(f'{k}="{v}"' for k, v in elem.attrib.items())
    print('  ' * indent + f'<{tag} {attrs}>'.rstrip())
    for child in elem:
        show(child, indent + 1)

slide = prs.slides[0]
for shape in slide.shapes:
    if shape.name == 'title text':
        txbody = shape._element.find('.//p:txBody', nsmap)
        print('--- title text txBody ---')
        show(txbody)
        break
