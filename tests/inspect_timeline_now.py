"""检查 timeline PPT 的文本与填充属性。"""
from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent / "src"))

from pptx import Presentation as PptxPresentation
from ppt_transfor.parser.presentation import parse_presentation


def main():
    pptx_path = Path("input/content_page__component_timeline.pptx")
    print("=== RAW PPTX ===")
    prs = PptxPresentation(str(pptx_path))
    for si, slide in enumerate(prs.slides):
        print(f"Slide {si}:")
        for shape in slide.shapes:
            print(f"  Shape: {shape.name}, type={shape.shape_type}")
            if hasattr(shape, "fill"):
                try:
                    ft = shape.fill.type
                    print(f"    fill.type={ft}")
                except Exception as e:
                    print(f"    fill.error={e}")
            if shape.has_text_frame:
                for pi, para in enumerate(shape.text_frame.paragraphs):
                    for ri, run in enumerate(para.runs):
                        text = run.text.strip()
                        if text:
                            try:
                                c = run.font.color
                                ct = c.type
                                rgb = getattr(c, "rgb", None)
                                print(f"    text[{pi},{ri}]: '{text[:40]}' color_type={ct} color={rgb}")
                            except Exception as e:
                                print(f"    text[{pi},{ri}]: '{text[:40]}' color_error={e}")

    print("\n=== PARSED MODEL ===")
    model = parse_presentation(str(pptx_path))
    for si, slide in enumerate(model.slides):
        print(f"Slide {si} bg={slide.background}")
        for shape in slide.shapes:
            name = getattr(shape, "name", None)
            print(f"  Shape: {name} type={type(shape).__name__}")
            if hasattr(shape, "fill"):
                print(f"    fill={shape.fill}")
            if hasattr(shape, "text"):
                for pi, para in enumerate(shape.text.paragraphs):
                    for ri, run in enumerate(para.runs):
                        if run.text.strip():
                            print(f"    text[{pi},{ri}]: '{run.text[:40]}' color={run.font.color}")


if __name__ == "__main__":
    main()
