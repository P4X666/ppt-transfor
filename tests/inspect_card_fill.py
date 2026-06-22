"""检查 card PPT 中特定文本框的 fill 类型。"""
from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent / "src"))

from pptx import Presentation as PptxPresentation


def main():
    prs = PptxPresentation("input/conent_page_component_card.pptx")
    slide = prs.slides[1]  # slide2 (0-indexed)
    for shape in slide.shapes:
        if shape.shape_type == 6:  # GROUP
            print(f"Group: {shape.name}")
            for child in shape.shapes:
                print(f"  Child: {child.name} type={child.shape_type}")
                if child.has_text_frame:
                    for para in child.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                print(f"    text: '{run.text[:30]}'")
                try:
                    ft = child.fill.type
                    print(f"    fill.type={ft} name={ft.name if hasattr(ft, 'name') else 'N/A'}")
                except Exception as e:
                    print(f"    fill.error={e}")
                # check XML
                try:
                    spPr = child._element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}spPr")
                    if spPr is not None:
                        print(f"    spPr children: {[c.tag.split('}')[-1] for c in spPr[:5]]}")
                except Exception as e:
                    print(f"    xml.error={e}")


if __name__ == "__main__":
    main()
